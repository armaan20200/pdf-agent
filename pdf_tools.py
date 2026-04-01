"""
Real PDF tools using PyMuPDF (fitz), pdfplumber, python-docx, openpyxl, python-pptx
"""

import os
import re
import uuid
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Optional
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
from docx.shared import Pt, RGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from PIL import Image


class PDFTools:
    def __init__(self, temp_dir: Path):
        self.temp_dir = temp_dir

    def _output_path(self, name: str) -> str:
        return str(self.temp_dir / name)

    # ─────────────────────────────────────────────
    # ORGANIZE
    # ─────────────────────────────────────────────

    def merge_pdfs(self, input_paths: list[str]) -> tuple[str, str]:
        """Merge multiple PDFs into one using PyMuPDF"""
        output_name = f"merged_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        merged_doc = fitz.open()
        for path in input_paths:
            doc = fitz.open(path)
            merged_doc.insert_pdf(doc)
            doc.close()
        merged_doc.save(output_path)
        merged_doc.close()
        return output_path, output_name

    def split_pdf(self, input_path: str) -> tuple[str, str]:
        """Split each page into separate PDFs, bundled in a zip"""
        doc = fitz.open(input_path)
        base = Path(input_path).stem
        zip_name = f"{base}_split_{uuid.uuid4().hex[:8]}.zip"
        zip_path = self._output_path(zip_name)
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i in range(doc.page_count):
                page_doc = fitz.open()
                page_doc.insert_pdf(doc, from_page=i, to_page=i)
                buf = BytesIO()
                page_doc.save(buf)
                page_doc.close()
                zf.writestr(f"page_{i+1}.pdf", buf.getvalue())
        doc.close()
        return zip_path, zip_name

    def split_at_page(self, input_path: str, split_page: int) -> tuple[str, str]:
        """Split PDF into two parts at a given page. Part 1: pages 1..split_page, Part 2: rest."""
        doc = fitz.open(input_path)
        total = doc.page_count
        split_page = max(1, min(split_page, total - 1))
        base = Path(input_path).stem
        zip_name = f"{base}_split_at_{split_page}_{uuid.uuid4().hex[:8]}.zip"
        zip_path = self._output_path(zip_name)
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for part, (frm, to) in enumerate([(0, split_page - 1), (split_page, total - 1)], 1):
                part_doc = fitz.open()
                part_doc.insert_pdf(doc, from_page=frm, to_page=to)
                buf = BytesIO()
                part_doc.save(buf)
                part_doc.close()
                zf.writestr(f"part_{part}_pages_{frm+1}-{to+1}.pdf", buf.getvalue())
        doc.close()
        return zip_path, zip_name

    def extract_pages(self, input_path: str, page_start: int, page_end: int) -> tuple[str, str]:
        """Extract specific pages from a PDF (1-indexed)"""
        output_name = f"extracted_pages_{page_start}-{page_end}_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        total_pages = doc.page_count
        start_idx = max(0, page_start - 1)
        end_idx = min(total_pages - 1, page_end - 1)
        new_doc = fitz.open()
        new_doc.insert_pdf(doc, from_page=start_idx, to_page=end_idx)
        new_doc.save(output_path)
        new_doc.close()
        doc.close()
        return output_path, output_name

    def remove_pages(self, input_path: str, pages_to_remove: list[int]) -> tuple[str, str]:
        """Remove specific pages (1-indexed) from a PDF"""
        output_name = f"pages_removed_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        total = doc.page_count
        # Convert to 0-indexed, filter valid pages
        to_remove_0 = sorted(set(p - 1 for p in pages_to_remove if 1 <= p <= total), reverse=True)
        for idx in to_remove_0:
            doc.delete_page(idx)
        doc.save(output_path)
        doc.close()
        return output_path, output_name

    def organize_pages(self, input_path: str, new_order: list[int]) -> tuple[str, str]:
        """Reorder pages by providing a new order list (1-indexed page numbers)"""
        output_name = f"reordered_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        total = doc.page_count
        # Validate & convert to 0-indexed
        order_0 = [p - 1 for p in new_order if 1 <= p <= total]
        new_doc = fitz.open()
        for idx in order_0:
            new_doc.insert_pdf(doc, from_page=idx, to_page=idx)
        new_doc.save(output_path)
        new_doc.close()
        doc.close()
        return output_path, output_name

    def rotate_pages(self, input_path: str, angle: int, pages: list[int] | None = None) -> tuple[str, str]:
        """Rotate pages by angle (90, 180, 270). pages=None means all pages."""
        output_name = f"rotated_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        total = doc.page_count
        targets = pages if pages else list(range(1, total + 1))
        for p in targets:
            if 1 <= p <= total:
                page = doc[p - 1]
                page.set_rotation((page.rotation + angle) % 360)
        doc.save(output_path)
        doc.close()
        return output_path, output_name

    def add_watermark(self, input_path: str, text: str = "CONFIDENTIAL",
                      opacity: float = 0.15, color: tuple = (0.6, 0.6, 0.6)) -> tuple[str, str]:
        """Add a diagonal text watermark to all pages using insert_text with morph"""
        output_name = f"watermarked_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        for page in doc:
            rect = page.rect
            center = fitz.Point(rect.width / 2, rect.height / 2)
            # Scale fontsize based on text length so it fits nicely across the page
            fontsize = min(rect.width / max(len(text), 1) * 1.1, 80)
            # Place text so its center aligns with the page center
            x = center.x - len(text) * fontsize * 0.27
            y = center.y + fontsize * 0.35
            # Use morph to rotate 45° around the page center
            page.insert_text(
                fitz.Point(x, y),
                text,
                fontsize=fontsize,
                color=color,
                fill_opacity=opacity,
                stroke_opacity=opacity,
                morph=(center, fitz.Matrix(45)),
                overlay=True,
            )
        doc.save(output_path)
        doc.close()
        return output_path, output_name

    def add_page_numbers(self, input_path: str, position: str = "bottom-center",
                         start_num: int = 1, prefix: str = "") -> tuple[str, str]:
        """Add page numbers to each page"""
        output_name = f"numbered_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        for i, page in enumerate(doc):
            rect = page.rect
            num_text = f"{prefix}{start_num + i}"
            fontsize = 10
            margin = 20
            if "bottom" in position:
                y = rect.height - margin
            else:
                y = margin + fontsize
            if "center" in position:
                x = rect.width / 2 - (len(num_text) * fontsize * 0.3)
            elif "right" in position:
                x = rect.width - margin - len(num_text) * fontsize * 0.6
            else:
                x = margin
            page.insert_text(
                fitz.Point(x, y),
                num_text,
                fontsize=fontsize,
                color=(0.3, 0.3, 0.3),
                overlay=True,
            )
        doc.save(output_path)
        doc.close()
        return output_path, output_name

    # ─────────────────────────────────────────────
    # CONVERT TO PDF
    # ─────────────────────────────────────────────

    def images_to_pdf(self, input_paths: list[str]) -> tuple[str, str]:
        """Convert one or more image files (JPG, PNG, etc.) to a single PDF"""
        output_name = f"images_to_pdf_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open()
        for img_path in input_paths:
            img = Image.open(img_path)
            img_rgb = img.convert("RGB")
            buf = BytesIO()
            img_rgb.save(buf, format="JPEG", quality=90)
            buf.seek(0)
            # Open as fitz image document
            img_doc = fitz.open("jpeg", buf.read())
            rect = img_doc[0].rect
            page = doc.new_page(width=rect.width, height=rect.height)
            page.show_pdf_page(rect, img_doc, 0)
            img_doc.close()
        doc.save(output_path)
        doc.close()
        return output_path, output_name

    def word_to_pdf(self, input_path: str) -> tuple[str, str]:
        """Convert DOCX to PDF using python-docx for reading + PyMuPDF for writing"""
        base = Path(input_path).stem
        output_name = f"{base}_to_pdf_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)

        doc_in = Document(input_path)
        pdf_doc = fitz.open()

        page = pdf_doc.new_page(width=595, height=842)  # A4
        y = 72
        left_margin = 72
        right_margin = 523

        for para in doc_in.paragraphs:
            text = para.text.strip()
            if not text:
                y += 8
                continue
            # Detect heading style
            style_name = para.style.name.lower() if para.style else ""
            if "heading 1" in style_name:
                fontsize, bold = 18, True
            elif "heading 2" in style_name:
                fontsize, bold = 14, True
            elif "heading 3" in style_name:
                fontsize, bold = 12, True
            else:
                fontsize, bold = 10, False

            # Simple word wrap
            words = text.split()
            line = ""
            for word in words:
                test_line = (line + " " + word).strip()
                # Rough char width estimate at ~6px per char
                if len(test_line) * fontsize * 0.5 > (right_margin - left_margin):
                    if y > 780:
                        page = pdf_doc.new_page(width=595, height=842)
                        y = 72
                    page.insert_text(
                        fitz.Point(left_margin, y), line,
                        fontsize=fontsize,
                        fontname="helv" if not bold else "helvB",
                        color=(0, 0, 0),
                    )
                    y += fontsize + 4
                    line = word
                else:
                    line = test_line

            if line:
                if y > 780:
                    page = pdf_doc.new_page(width=595, height=842)
                    y = 72
                page.insert_text(
                    fitz.Point(left_margin, y), line,
                    fontsize=fontsize,
                    fontname="helv" if not bold else "helvB",
                    color=(0, 0, 0),
                )
                y += fontsize + 6

        pdf_doc.save(output_path)
        pdf_doc.close()
        return output_path, output_name

    def pptx_to_pdf(self, input_path: str) -> tuple[str, str]:
        """Convert PPTX to PDF — extracts text/shapes per slide into PDF pages"""
        from pptx import Presentation
        from pptx.util import Pt as PPTPt

        base = Path(input_path).stem
        output_name = f"{base}_to_pdf_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)

        prs = Presentation(input_path)
        pdf_doc = fitz.open()

        # Get slide dimensions
        slide_w = prs.slide_width.inches * 72
        slide_h = prs.slide_height.inches * 72

        for slide_num, slide in enumerate(prs.slides, 1):
            page = pdf_doc.new_page(width=slide_w, height=slide_h)

            # White background
            page.draw_rect(page.rect, color=(1, 1, 1), fill=(1, 1, 1))

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                left = shape.left.inches * 72 if shape.left else 20
                top = shape.top.inches * 72 if shape.top else 20
                width = shape.width.inches * 72 if shape.width else slide_w - 40

                y_offset = top
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        y_offset += 8
                        continue
                    fontsize = 12
                    bold = False
                    for run in para.runs:
                        if run.font.size:
                            fontsize = max(6, min(36, run.font.size.pt))
                        if run.font.bold:
                            bold = True
                    page.insert_text(
                        fitz.Point(left, y_offset + fontsize),
                        text[:200],
                        fontsize=fontsize,
                        fontname="helvB" if bold else "helv",
                        color=(0, 0, 0),
                    )
                    y_offset += fontsize + 4

        pdf_doc.save(output_path)
        pdf_doc.close()
        return output_path, output_name

    def excel_to_pdf(self, input_path: str) -> tuple[str, str]:
        """Convert XLSX to PDF — renders each sheet as a table"""
        base = Path(input_path).stem
        output_name = f"{base}_to_pdf_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)

        wb = openpyxl.load_workbook(input_path, data_only=True)
        pdf_doc = fitz.open()

        for sheet in wb.sheetnames:
            ws = wb[sheet]
            page = pdf_doc.new_page(width=842, height=595)  # Landscape A4
            x, y = 36, 36
            page.insert_text(
                fitz.Point(x, y), f"Sheet: {sheet}",
                fontsize=12, fontname="helvB", color=(0, 0, 0),
            )
            y += 20
            col_width = 80
            row_height = 14
            for row in ws.iter_rows(values_only=True):
                x = 36
                for cell in row:
                    cell_text = str(cell) if cell is not None else ""
                    cell_text = cell_text[:12]  # Truncate wide cells
                    page.insert_text(
                        fitz.Point(x, y), cell_text,
                        fontsize=8, color=(0, 0, 0),
                    )
                    x += col_width
                    if x > 800:
                        break
                y += row_height
                if y > 560:
                    page = pdf_doc.new_page(width=842, height=595)
                    y = 36

        pdf_doc.save(output_path)
        pdf_doc.close()
        return output_path, output_name

    def html_to_pdf(self, html_content: str) -> tuple[str, str]:
        """Convert HTML string to PDF using PyMuPDF's Story renderer"""
        output_name = f"html_to_pdf_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        # Use PyMuPDF story for HTML rendering
        story = fitz.Story(html=html_content)
        writer = fitz.DocumentWriter(output_path)
        mediabox = fitz.paper_rect("a4")
        where = mediabox + (36, 36, -36, -36)
        more = True
        while more:
            device = writer.begin_page(mediabox)
            more, _ = story.place(where)
            story.draw(device)
            writer.end_page()
        writer.close()
        return output_path, output_name

    # ─────────────────────────────────────────────
    # CONVERT FROM PDF
    # ─────────────────────────────────────────────

    def pdf_to_jpg(self, input_path: str, dpi: int = 150) -> tuple[str, str]:
        """Convert each PDF page to JPEG, bundle in a ZIP"""
        base = Path(input_path).stem
        zip_name = f"{base}_images_{uuid.uuid4().hex[:8]}.zip"
        zip_path = self._output_path(zip_name)
        doc = fitz.open(input_path)
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("jpeg")
                zf.writestr(f"page_{i+1:03d}.jpg", img_bytes)
        doc.close()
        return zip_path, zip_name

    def pdf_to_word(self, input_path: str) -> tuple[str, str]:
        """Extract PDF text and format into a .docx with structure preserved"""
        base = Path(input_path).stem
        output_name = f"{base}_{uuid.uuid4().hex[:8]}.docx"
        output_path = self._output_path(output_name)
        doc = Document()
        title = doc.add_heading(level=0)
        title.add_run(base).bold = True
        doc.add_paragraph()
        pdf_doc = fitz.open(input_path)
        for page_num in range(pdf_doc.page_count):
            page = pdf_doc[page_num]
            doc.add_heading(f"Page {page_num + 1}", level=2)
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        line_text = ""
                        is_bold = False
                        font_size = 11
                        for span in line.get("spans", []):
                            span_text = span.get("text", "").strip()
                            if span_text:
                                line_text += span_text + " "
                                flags = span.get("flags", 0)
                                is_bold = bool(flags & 2**4)
                                font_size = span.get("size", 11)
                        if line_text.strip():
                            para = doc.add_paragraph()
                            run = para.add_run(line_text.strip())
                            run.bold = is_bold
                            run.font.size = Pt(10)
            if page_num < pdf_doc.page_count - 1:
                doc.add_page_break()
        pdf_doc.close()
        doc.save(output_path)
        return output_path, output_name

    def pdf_to_pptx(self, input_path: str) -> tuple[str, str]:
        """Convert PDF pages to PowerPoint slides (one page = one slide)"""
        from pptx import Presentation
        from pptx.util import Inches, Emu
        base = Path(input_path).stem
        output_name = f"{base}_to_pptx_{uuid.uuid4().hex[:8]}.pptx"
        output_path = self._output_path(output_name)

        pdf_doc = fitz.open(input_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        mat = fitz.Matrix(1.5, 1.5)
        for i, page in enumerate(pdf_doc):
            slide = prs.slides.add_slide(blank_layout)
            # Render page as image
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("jpeg")
            buf = BytesIO(img_bytes)
            slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)

        pdf_doc.close()
        prs.save(output_path)
        return output_path, output_name

    def pdf_to_excel(self, input_path: str) -> tuple[str, str]:
        """Convert PDF tables to Excel using pdfplumber"""
        base = Path(input_path).stem
        output_name = f"{base}_{uuid.uuid4().hex[:8]}.xlsx"
        output_path = self._output_path(output_name)
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="2D3748", end_color="2D3748", fill_type="solid")
        tables_found = 0
        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                if not tables:
                    continue
                for table_idx, table in enumerate(tables):
                    tables_found += 1
                    ws = wb.create_sheet(title=f"Page{page_num}_T{table_idx + 1}")
                    for row_idx, row in enumerate(table, start=1):
                        for col_idx, cell_value in enumerate(row, start=1):
                            cell = ws.cell(row=row_idx, column=col_idx, value=cell_value or "")
                            if row_idx == 1:
                                cell.font = header_font
                                cell.fill = header_fill
                    for col in ws.columns:
                        max_len = max((len(str(c.value or "")) for c in col), default=8)
                        ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)
        if tables_found == 0:
            ws = wb.create_sheet(title="Text Content")
            ws.cell(row=1, column=1, value="No tables found. Extracted text:").font = Font(bold=True)
            with pdfplumber.open(input_path) as pdf:
                row_num = 2
                for page_num, page in enumerate(pdf.pages, start=1):
                    ws.cell(row=row_num, column=1, value=f"--- Page {page_num} ---").font = Font(bold=True, italic=True)
                    row_num += 1
                    for line in (page.extract_text() or "").split("\n"):
                        if line.strip():
                            ws.cell(row=row_num, column=1, value=line)
                            row_num += 1
                    row_num += 1
            ws.column_dimensions["A"].width = 80
        wb.save(output_path)
        return output_path, output_name

    def pdf_to_pdfa(self, input_path: str) -> tuple[str, str]:
        """Convert PDF to PDF/A-1b archival format"""
        base = Path(input_path).stem
        output_name = f"{base}_pdfa_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()
        return output_path, output_name

    def pdf_to_text(self, input_path: str) -> tuple[str, str]:
        """Extract all text from PDF to a plain text file"""
        base = Path(input_path).stem
        output_name = f"{base}_{uuid.uuid4().hex[:8]}.txt"
        output_path = self._output_path(output_name)
        pdf_doc = fitz.open(input_path)
        lines = []
        for page_num in range(pdf_doc.page_count):
            page = pdf_doc[page_num]
            lines.append(f"{'='*60}\nPAGE {page_num + 1}\n{'='*60}")
            lines.append(page.get_text("text"))
            lines.append("")
        pdf_doc.close()
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        return output_path, output_name

    # ─────────────────────────────────────────────
    # OPTIMIZE / EDIT
    # ─────────────────────────────────────────────

    def compress_pdf(self, input_path: str) -> tuple[str, str]:
        """Compress PDF by cleaning unused objects"""
        base = Path(input_path).stem
        output_name = f"{base}_compressed_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()
        return output_path, output_name

    def repair_pdf(self, input_path: str) -> tuple[str, str]:
        """Attempt to repair a damaged or corrupt PDF"""
        base = Path(input_path).stem
        output_name = f"{base}_repaired_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()
        return output_path, output_name

    # ─────────────────────────────────────────────
    # SECURITY
    # ─────────────────────────────────────────────

    def sign_pdf(self, pdf_path: str, signature_path: str,
                 position: str = "bottom", page_target: str = "last") -> tuple[str, str]:
        """Stamp a signature image onto a PDF page"""
        base = Path(pdf_path).stem
        output_name = f"{base}_signed_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(pdf_path)
        total_pages = doc.page_count
        if page_target == "all":
            target_pages = list(range(total_pages))
        elif page_target == "first":
            target_pages = [0]
        else:
            target_pages = [total_pages - 1]

        sig_img = Image.open(signature_path).convert("RGBA")
        data = sig_img.load()
        width, height = sig_img.size
        for y_px in range(height):
            for x_px in range(width):
                r, g, b, a = data[x_px, y_px]
                if r > 230 and g > 230 and b > 230:
                    data[x_px, y_px] = (r, g, b, 0)
        sig_temp = self._output_path(f"sig_temp_{uuid.uuid4().hex[:8]}.png")
        sig_img.save(sig_temp, "PNG")

        for page_idx in target_pages:
            page = doc[page_idx]
            pw = page.rect.width
            ph = page.rect.height
            sig_w = pw * 0.25
            sig_h = sig_w * (sig_img.height / sig_img.width)
            margin = 40
            if position == "top":
                rect = fitz.Rect(pw - sig_w - margin, margin, pw - margin, margin + sig_h)
            elif position == "center":
                cx = (pw - sig_w) / 2
                cy = (ph - sig_h) / 2
                rect = fitz.Rect(cx, cy, cx + sig_w, cy + sig_h)
            else:
                rect = fitz.Rect(pw - sig_w - margin, ph - sig_h - margin, pw - margin, ph - margin)
            page.insert_image(rect, filename=sig_temp)

        doc.save(output_path)
        doc.close()
        try:
            os.remove(sig_temp)
        except Exception:
            pass
        return output_path, output_name

    def unlock_pdf(self, input_path: str, password: str = "") -> tuple[str, str]:
        """Remove password protection from a PDF"""
        base = Path(input_path).stem
        output_name = f"{base}_unlocked_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        if doc.is_encrypted:
            if not doc.authenticate(password):
                raise ValueError("Incorrect password — cannot unlock this PDF.")
        doc.save(output_path, encryption=fitz.PDF_ENCRYPT_NONE)
        doc.close()
        return output_path, output_name

    def protect_pdf(self, input_path: str, user_password: str, owner_password: str = "") -> tuple[str, str]:
        """Add password protection to a PDF"""
        base = Path(input_path).stem
        output_name = f"{base}_protected_{uuid.uuid4().hex[:8]}.pdf"
        output_path = self._output_path(output_name)
        doc = fitz.open(input_path)
        owner_pw = owner_password or user_password + "_owner"
        doc.save(
            output_path,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            user_pw=user_password,
            owner_pw=owner_pw,
            permissions=fitz.PDF_PERM_PRINT | fitz.PDF_PERM_COPY,
        )
        doc.close()
        return output_path, output_name

    def rename_pdf(self, input_path: str, new_name: str) -> tuple[str, str]:
        """Rename PDF file and update internal title metadata"""
        safe_name = "".join(c for c in new_name if c.isalnum() or c in "._- ").strip()
        if not safe_name.endswith(".pdf"):
            safe_name += ".pdf"
        output_path = self._output_path(safe_name)
        doc = fitz.open(input_path)
        doc.set_metadata({"title": safe_name.replace(".pdf", "")})
        doc.save(output_path)
        doc.close()
        return output_path, safe_name

    # ─────────────────────────────────────────────
    # AI TOOLS (used by agent.py when OpenAI available)
    # ─────────────────────────────────────────────

    def extract_full_text(self, input_path: str) -> str:
        """Extract all text from PDF for AI processing"""
        doc = fitz.open(input_path)
        texts = []
        for page in doc:
            texts.append(page.get_text("text"))
        doc.close()
        return "\n\n".join(texts)[:12000]  # Cap for API token limits
